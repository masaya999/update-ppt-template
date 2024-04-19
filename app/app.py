
from flask import Flask, render_template, request, send_from_directory
import pandas as pd
from os import listdir, path
from pptx import Presentation
from reportlab.pdfgen import canvas

app = Flask(__name__)

csv_path = '../data/extracted_data/saved_contents.csv'
pptx_base_path = '../data/pptx/'
# CSVファイルを開いて読み込む
df = pd.read_csv(csv_path, encoding='cp932', index_col=0)
file_list = sorted(set(df['file_name']))
client_info = {}
proposal_info = {}
selected_file = None

def load_contents(file_name):
    df_tmp = df[df['file_name']==file_name]
    df_client = df_tmp[df_tmp['info_level']=='client']
    for index, row in df_client.iterrows():
        client_info[row['label']] = row['content']
    df_proposal = df_tmp[df_tmp['info_level']=='proposal']
    for index, row in df_proposal.iterrows():
        proposal_info[row['label']] = row['content']
    #return client_info, proposal_info

def convert_pptx_to_pdf(pptx_filename, pdf_filename):
    prs = Presentation(pptx_filename)
    c = canvas.Canvas(pdf_filename)
    for slide in prs.slides:
        c.drawString(100, 750, "Slide Title Placeholder")  # 簡単なプレースホルダー
        c.showPage()
    c.save()
    
@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        selected_file = request.form.get('file_selector')
        load_contents(selected_file)
        return render_template('index.html', files=file_list, selected_file=selected_file, client_info=client_info, proposal_info=proposal_info)
    return render_template('index.html', files=file_list, selected_file=None, client_info=None, proposal_info=None)

@app.route('/files/<filename>')
def file(filename):
    # ファイルがPDFでなければ変換
    base, ext = path.splitext(filename)
    if ext.lower() == ".pptx":
        pdf_filename = base + ".pdf"
        convert_pptx_to_pdf(filename, pdf_filename)
        filename = pdf_filename
    return send_from_directory('.', filename)

@app.route('/update/client', methods=['POST'])###temp!!
def update_client():
    # テキストボックスからの新しいデータを受け取り、dictを更新
    for key in client_info.keys():
        client_info[key] = request.form.get(key)
    return render_template('index.html', files=file_list, selected_file=selected_file, client_info=client_info, proposal_info=proposal_info)

@app.route('/update/proposal', methods=['POST'])###temp!!
def update_proposal():
    # テキストボックスからの新しいデータを受け取り、dictを更新
    for key in proposal_info.keys():
        proposal_info[key] = request.form.get(key)
    return render_template('index.html', files=file_list, selected_file=selected_file, client_info=client_info, proposal_info=proposal_info)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0')
"""
@app.route('/')
def index():
    
    return render_template('index.html', file_option=file_list)


    #custom_text = copy.deepcopy(core_text)
    #custom_text = {key: value for key, value in core_text.items()}
    print(core_text)
    print(custom_text)
    return render_template('index.html', data=custom_text)

@app.route('/update', methods=['POST'])
def update():
    # テキストボックスからの新しいデータを受け取り、dictを更新
    for key in custom_text.keys():
        custom_text[key] = request.get(key)



@app.route('/update', methods=['POST'])
def update():
    # テキストボックスからの新しいデータを受け取り、dictを更新
    for key in custom_text.keys():
        custom_text[key] = request.get(key)
    
    # PowerPointファイルを更新する関数を呼び出す
    update_pptx_with_custom_text()
    
    return render_template('index.html', data=custom_text, filename='customized_'+pptx_file_path)

def update_pptx_with_custom_text():
    prs = Presentation(pptx_file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in core_text.items():
                            if value in run.text and custom_text[key] != value:
                                run.text = run.text.replace(value, custom_text[key])
    prs.save('customized_'+pptx_file_path)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0')


from flask import Flask, render_template, request, send_from_directory
import pandas as pd
from os import listdir, path
from pptx import Presentation
from reportlab.pdfgen import canvas

app = Flask(__name__)

def load_csv_data(filename):
    return pd.read_csv(filename)

def convert_pptx_to_pdf(pptx_filename, pdf_filename):
    prs = Presentation(pptx_filename)
    c = canvas.Canvas(pdf_filename)
    for slide in prs.slides:
        c.drawString(100, 750, "Slide Title Placeholder")  # 簡単なプレースホルダー
        c.showPage()
    c.save()

@app.route('/')
def home():
    csv_data = load_csv_data('data.csv')
    pptx_files = [f for f in listdir('.') if f.endswith('.pptx')]
    return render_template('index.html', data=csv_data, pptx_files=pptx_files)

@app.route('/files/<filename>')
def file(filename):
    # ファイルがPDFでなければ変換
    base, ext = path.splitext(filename)
    if ext.lower() == ".pptx":
        pdf_filename = base + ".pdf"
        convert_pptx_to_pdf(filename, pdf_filename)
        filename = pdf_filename
    return send_from_directory('.', filename)

if __name__ == '__main__':
    app.run(debug=True)
"""
